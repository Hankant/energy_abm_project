from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


OUTPUT = "c:/Users/69596/.gemini/antigravity/scratch/energy_abm_project/docs/planning/research_logic_presentation.pptx"


def set_run_font(run, size=20, bold=False, color=(34, 34, 34), name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, title_text, subtitle_text=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    set_run_font(r, size=26, bold=True, color=(15, 52, 96))

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(11), Inches(0.45))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle_text
        set_run_font(r2, size=11, color=(90, 90, 90))

    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.55), Inches(11.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(15, 52, 96)
    line.line.color.rgb = RGBColor(15, 52, 96)


def add_bullets(slide, items, left=0.95, top=1.9, width=11.0, height=5.0, level0=21, level1=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        p.bullet = True
        r = p.add_run()
        r.text = text
        set_run_font(r, size=level0 if level == 0 else level1)


def add_center_text(slide, text, left, top, width, height, size=18, bold=False, color=(34, 34, 34)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, bold=bold, color=color)


def add_logic_boxes(slide, labels):
    left = 0.55
    top = 2.35
    width = 2.15
    height = 1.2
    gap = 0.2
    colors = [
        (224, 236, 248),
        (215, 231, 210),
        (255, 239, 204),
        (242, 220, 219),
        (221, 217, 238),
    ]
    for i, label in enumerate(labels):
        x = Inches(left + i * (width + gap))
        shape = slide.shapes.add_shape(1, x, Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors[i % len(colors)])
        shape.line.color.rgb = RGBColor(120, 120, 120)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_font(r, size=18, bold=True, color=(45, 45, 45))
        if i < len(labels) - 1:
            add_center_text(slide, "→", left + i * (width + gap) + width + 0.02, top + 0.32, 0.16, 0.35, size=24, bold=True, color=(60, 60, 60))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    slide = prs.slides.add_slide(blank)
    add_center_text(slide, "多维政策组合如何在极端气候下促进家庭能源结构低碳转型", 0.9, 1.1, 11.5, 0.9, size=26, bold=True, color=(15, 52, 96))
    add_center_text(slide, "一项基于 LLM-Agents 的仿真研究", 1.8, 2.0, 9.7, 0.5, size=18, color=(70, 70, 70))
    add_center_text(slide, "开题汇报", 5.4, 2.65, 2.3, 0.45, size=18, bold=True, color=(110, 70, 20))
    add_logic_boxes(slide, ["极端气候", "政策组合", "家庭异质性", "社会互动", "转型结果"])

    # 2 Background
    slide = prs.slides.add_slide(blank)
    add_title(slide, "1. 研究背景：问题从哪里来")
    add_bullets(slide, [
        "极端高温、寒潮和持续异常天气正在显著改变家庭制冷、采暖和日常用电需求，并推高电力系统峰值负荷。",
        "家庭能源低碳转型并不是单纯的技术替换问题，而是一个同时受收入、价格、基础设施、风险认知和社会互动影响的复杂行为过程。",
        "现实中单一政策工具往往难以同时实现节能减排、削峰降负和公平改善，这使“政策组合”成为更有解释力的研究视角。",
        "因此，本研究关注的不是“某一种政策是否有效”，而是“在极端气候背景下，什么样的政策组合对什么样的家庭更有效且更公平”。",
    ])

    # 3 Literature gap
    slide = prs.slides.add_slide(blank)
    add_title(slide, "2. 文献缺口：现有研究缺什么")
    add_bullets(slide, [
        "已有气候与能源研究较充分解释了极端天气如何推高居民能源需求，但较少进一步讨论家庭低碳转型的行为路径。",
        "已有绿色消费和支付意愿研究能够识别个体偏好，却较难解释这些偏好在动态环境中如何累积为群体层面的转型结果。",
        "已有政策研究已经从单一工具走向政策组合，但更多聚焦宏观系统、产业技术和总体减排路径，对家庭层面的机制分析仍然不足。",
        "已有公平研究强调收入和区域差异的重要性，但缺乏一个能同时把效率、行为响应和分配结果放在一起分析的框架。",
    ])

    # 4 Core question
    slide = prs.slides.add_slide(blank)
    add_title(slide, "3. 核心研究问题")
    add_center_text(slide, "在极端气候持续冲击下，", 1.2, 2.0, 11.0, 0.5, size=24, bold=True, color=(15, 52, 96))
    add_center_text(slide, "多维政策组合如何通过影响异质性家庭的决策与社会扩散，", 1.0, 2.7, 11.3, 0.5, size=24, bold=True, color=(15, 52, 96))
    add_center_text(slide, "塑造家庭能源低碳转型的效率与公平结果？", 1.7, 3.4, 10.0, 0.5, size=24, bold=True, color=(15, 52, 96))
    add_bullets(slide, [
        "关键词：极端气候、政策组合、家庭异质性、社会互动、低碳采用、公平结果"
    ], top=4.55, height=1.0, level0=18)

    # 5 Logic chain
    slide = prs.slides.add_slide(blank)
    add_title(slide, "4. 研究内在逻辑")
    add_logic_boxes(slide, [
        "气候冲击\n改变需求与成本压力",
        "政策组合\n塑造激励与约束",
        "异质家庭\n产生差异化响应",
        "社会互动\n推动扩散或分化",
        "系统结果\n表现为效率与公平"
    ])
    add_bullets(slide, [
        "逻辑主线不是“政策直接产生结果”，而是“政策通过家庭决策机制和社会传播机制间接塑造结果”。",
        "因此，研究重点必须同时包含行为异质性和互动过程，而不能只停留在静态均值比较。"
    ], top=4.4, height=1.8, level0=18)

    # 6 Framework
    slide = prs.slides.add_slide(blank)
    add_title(slide, "5. 分析框架")
    add_bullets(slide, [
        "外生冲击层：极端气候通过高温、寒潮和持续异常天气改变家庭能源需求与舒适性约束。",
        "政策工具层：经济激励、基础设施支持和规制约束共同影响家庭对节电和设备替换的收益判断。",
        "微观决策层：不同家庭因收入、住房条件、设备保有量、价格敏感度和低碳偏好不同，形成差异化响应。",
        "社会扩散层：邻里交流、观察学习和社会规范改变家庭的采用阈值与行为更新速度。",
        "结果评价层：最终体现在总用电量、峰值负荷、低碳产品采用率、支付意愿变化和社会公平指数上。",
    ])

    # 7 Why survey + ABM
    slide = prs.slides.add_slide(blank)
    add_title(slide, "6. 为什么要“问卷 + 情景实验 + LLM-Agent 仿真”")
    add_bullets(slide, [
        "问卷回答的是：家庭在特定情境下愿不愿意、愿意到什么程度、为什么愿意。",
        "情景实验回答的是：当极端天气和政策条件发生变化时，家庭的节电、购买和支付意愿如何变化。",
        "LLM-Agent 仿真回答的是：许多具有异质性偏好的家庭，在持续气候冲击和社会互动中反复决策时，系统可能如何演化。",
        "三者之间的关系是：问卷提供微观基础，情景实验提供政策响应参数，仿真负责把这些微观差异推演到群体层面。",
    ])

    # 8 Design
    slide = prs.slides.add_slide(blank)
    add_title(slide, "7. 研究设计")
    add_bullets(slide, [
        "研究对象：具有不同收入、住房、设备存量、价格敏感度和社会影响易感性的家庭。",
        "气候情景：以连续 7 天极热、连续 30 天极热等场景刻画极端气候冲击。",
        "政策维度：重点考察经济激励、基础设施支持和规制约束三类工具及其组合。",
        "核心结果：关注低碳产品额外支付意愿、节电强度、采用概率，以及峰值负荷和公平结果。",
    ])

    # 9 Hypotheses
    slide = prs.slides.add_slide(blank)
    add_title(slide, "8. 研究假设")
    add_bullets(slide, [
        "H1：极端气候持续时间越长，家庭的节电强度、低碳产品采用意愿和相关支付意愿越高。",
        "H2：经济激励型政策能够显著提高低碳产品采用概率，且对价格敏感和中低收入家庭更有效。",
        "H3：规制约束型政策更有助于抑制高耗电行为，但也可能带来更强的分配压力。",
        "H4：基础设施支持能够降低“想转型但无法转型”的物理约束，并放大政策效果。",
        "H5：社会互动会加速低碳行为扩散，并改变政策组合的整体效果排序。",
    ], top=1.9, height=4.9, level0=17)

    # 10 Contributions
    slide = prs.slides.add_slide(blank)
    add_title(slide, "9. 预期贡献")
    add_bullets(slide, [
        "理论上：把极端气候、政策组合、家庭异质性和社会扩散纳入同一研究框架。",
        "方法上：把问卷情景实验与 LLM-Agent 仿真结合起来，连接微观偏好识别与群体动态演化分析。",
        "实证上：比较不同政策组合对低碳采用、负荷优化和公平改善的相对效果。",
        "政策上：回答“什么样的政策组合、对什么样的家庭、更有效且更公平”。",
    ])

    # 11 Closing
    slide = prs.slides.add_slide(blank)
    add_title(slide, "10. 汇报总结")
    add_center_text(slide, "本研究关注的核心不是单一政策是否有效，", 1.2, 2.0, 11.0, 0.5, size=23, bold=True, color=(15, 52, 96))
    add_center_text(slide, "而是极端气候背景下，政策组合如何通过家庭决策与社会扩散", 0.8, 2.7, 11.7, 0.5, size=23, bold=True, color=(15, 52, 96))
    add_center_text(slide, "共同塑造家庭能源低碳转型的效率与公平路径。", 1.5, 3.4, 10.3, 0.5, size=23, bold=True, color=(15, 52, 96))
    add_bullets(slide, [
        "一句话概括：从“家庭怎么想”推进到“很多这样想的家庭互动后，系统会怎么变”。"
    ], top=4.65, height=1.0, level0=18)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
